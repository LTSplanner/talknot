"""docs/BUYBACK_SLIDES.md から TalKnot ブランドカラーの .pptx を生成する。

使い方:  python3 scripts/build_buyback_pptx.py
出力:    docs/BUYBACK_TalKnot.pptx
"""
from __future__ import annotations

import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# --- TalKnot ブランドカラー（ui/theme.py と対応）---
CORAL = RGBColor(0xFF, 0x6F, 0x61)
INDIGO = RGBColor(0x6C, 0x5C, 0xE7)
SUNNY = RGBColor(0xFF, 0xC3, 0x6B)
CREAM = RGBColor(0xFF, 0xF9, 0xF4)
INK = RGBColor(0x2D, 0x2A, 0x32)
MUTED = RGBColor(0x8C, 0x87, 0x94)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SRC = "docs/BUYBACK_SLIDES.md"
OUT = "docs/BUYBACK_TalKnot.pptx"


def parse_slides(path: str) -> list[tuple[str, list[str]]]:
    md = open(path, encoding="utf-8").read().split("## 付録")[0]
    slides: list[tuple[str, list[str]]] = []
    for block in md.split("\n---\n"):
        if "### スライド" not in block:
            continue
        title, bullets = None, []
        for line in block.splitlines():
            s = line.strip()
            if not s or s.startswith("### スライド"):
                continue
            m = re.match(r"\*\*(.+?)\*\*$", s)
            if m and title is None:
                title = m.group(1)
            elif s.startswith("- "):
                bullets.append(s[2:].strip().replace("**", ""))
        if title:
            slides.append((title, bullets))
    return slides


def _fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    return box, tf


def build() -> None:
    slides = parse_slides(SRC)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    W, H = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]

    for i, (title, bullets) in enumerate(slides):
        slide = prs.slides.add_slide(blank)

        # 背景
        bg = slide.shapes.add_shape(1, 0, 0, W, H)  # 1 = rectangle
        _fill(bg, INDIGO if i == 0 else CREAM)

        if i == 0:
            # 表紙
            band = slide.shapes.add_shape(1, 0, Inches(3.05), W, Inches(0.12))
            _fill(band, CORAL)
            _, tf = _textbox(slide, Inches(1), Inches(2.0), Inches(11.3), Inches(1.2))
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = "TalKnot（トークノット）🪢"
            r.font.size = Pt(48); r.font.bold = True; r.font.color.rgb = WHITE
            _, tf2 = _textbox(slide, Inches(1), Inches(3.4), Inches(11.3), Inches(2.2))
            for j, b in enumerate(bullets):
                p = tf2.paragraphs[0] if j == 0 else tf2.add_paragraph()
                p.alignment = PP_ALIGN.CENTER
                r = p.add_run(); r.text = b
                r.font.size = Pt(20); r.font.color.rgb = RGBColor(0xEC, 0xE9, 0xFB)
                p.space_after = Pt(6)
            continue

        # コンテンツ：上部アクセント帯
        top_band = slide.shapes.add_shape(1, 0, 0, W, Inches(0.18))
        _fill(top_band, CORAL)

        # タイトル
        _, ttf = _textbox(slide, Inches(0.8), Inches(0.55), Inches(11.7), Inches(1.0))
        p = ttf.paragraphs[0]
        r = p.add_run(); r.text = title
        r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = INDIGO

        # タイトル下のコーラル下線
        ul = slide.shapes.add_shape(1, Inches(0.85), Inches(1.55), Inches(2.6), Inches(0.06))
        _fill(ul, CORAL)

        # 本文（箇条書き）
        _, btf = _textbox(slide, Inches(0.95), Inches(1.9), Inches(11.4), Inches(4.9))
        btf.vertical_anchor = MSO_ANCHOR.TOP
        for j, b in enumerate(bullets):
            p = btf.paragraphs[0] if j == 0 else btf.add_paragraph()
            dot = p.add_run(); dot.text = "● "
            dot.font.size = Pt(14); dot.font.color.rgb = CORAL
            r = p.add_run(); r.text = b
            r.font.size = Pt(19); r.font.color.rgb = INK
            p.space_after = Pt(12); p.line_spacing = 1.1

        # フッター（左：ブランド / 右：ページ番号）
        _, ftf = _textbox(slide, Inches(0.8), Inches(6.95), Inches(6), Inches(0.4))
        fr = ftf.paragraphs[0].add_run(); fr.text = "TalKnot 🪢  ｜ システム買取制度 申請資料"
        fr.font.size = Pt(11); fr.font.color.rgb = MUTED
        _, ptf = _textbox(slide, Inches(11.8), Inches(6.95), Inches(1.2), Inches(0.4))
        pp = ptf.paragraphs[0]; pp.alignment = PP_ALIGN.RIGHT
        pr = pp.add_run(); pr.text = f"{i+1} / {len(slides)}"
        pr.font.size = Pt(11); pr.font.color.rgb = MUTED

    prs.save(OUT)
    print(f"✅ 生成: {OUT}  （{len(slides)}スライド・ブランドカラー版）")


if __name__ == "__main__":
    build()
